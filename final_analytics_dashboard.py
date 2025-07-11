# final_analytics_dashboard.py – Enhanced UI Version
from __future__ import annotations
import json, uuid, re
from datetime import datetime, timezone
from pathlib import Path

import streamlit as st
import pandas as pd
import chromadb
import altair as alt
from sentence_transformers import SentenceTransformer
import torch
import pytesseract

from faq_dashboard_with_creator import clean_and_tag
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ── Paths ──────────────────────────────────────────────────────────
BASE_DIR         = Path(__file__).parent
HISTORY_FILE     = BASE_DIR / "search_history.json"
QUERY_LOG_FILE   = BASE_DIR / "query_log.json"
FEEDBACK_FILE    = BASE_DIR / "feedback_log.json"
UPLOAD_LOG_FILE  = BASE_DIR / "upload_log.json"
FAQ_JSON_FILE    = BASE_DIR / "flocard_faq.json"
CHROMA_DIR       = BASE_DIR / "chroma_ui_storage"
COLL_NAME        = "week3_collection"

# ── Helpers ─────────────────────────────────────────────────────────
def now_iso() -> str:
    return datetime.now(timezone.utc).replace(tzinfo=None).isoformat(timespec="seconds")

def load_json(path: Path) -> list[dict]:
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text("utf-8") or "[]")
    except Exception:
        return []

def save_json(path: Path, data: list[dict]):
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")

def append_to_faq_json(new_entries: list[dict]):
    faq_data = load_json(FAQ_JSON_FILE)
    max_id = max((int(e.get("id", 0)) for e in faq_data), default=0)
    for i, entry in enumerate(new_entries, 1):
        entry_id = str(max_id + i)
        faq_data.append({
            "id": entry_id,
            "text": entry["text"],
            "metadata": entry["metadata"]
        })
    save_json(FAQ_JSON_FILE, faq_data)

def to_df(data: list[dict]) -> pd.DataFrame:
    df = pd.DataFrame(data)
    if df.empty:
        return df
    if "timestamp" in df.columns:
        dt = pd.to_datetime(df["timestamp"], errors="coerce", utc=True)
        df["timestamp"] = dt.dt.tz_localize(None)
        df = df.sort_values("timestamp", ascending=False)
    for col in df.columns:
        if df[col].apply(lambda x: isinstance(x, (list, dict))).any():
            df[col] = df[col].apply(
                lambda x: json.dumps(x, ensure_ascii=False, default=str)
                if isinstance(x, (list, dict)) else x
            )
    return df

# ── Load logs ───────────────────────────────────────────────────────
history   = load_json(HISTORY_FILE)
query_log = load_json(QUERY_LOG_FILE)
feedback  = load_json(FEEDBACK_FILE)
uploads   = load_json(UPLOAD_LOG_FILE)

# ── Ensure timestamps in feedback ───────────────────────────────────
for f in feedback:
    if "timestamp" not in f:
        f["timestamp"] = now_iso()

# ── Vector DB connection ────────────────────────────────────────────
embedder  = SentenceTransformer("all-MiniLM-L6-v2", device="cpu") #sentence transformer for embeddings
#all minilm-l6-v2 is a lightweight model suitable for many tasks and used here for generating embeddings
client    = chromadb.PersistentClient(path=str(CHROMA_DIR))
collection = client.get_or_create_collection(COLL_NAME)

# ── Page config ─────────────────────────────────────────────────────
st.set_page_config(page_title="Analytics Dashboard", layout="wide", page_icon="📊")
st.title("📊 Search Analytics Dashboard")

# Custom CSS for better UI
st.markdown("""
<style>
    .metric-card {
        background-color: #2E2E2E;
        border-radius: 10px;
        padding: 15px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    .stTabs [role="tablist"] {
        gap: 10px;
    }
    .stTabs [role="tab"] {
        border-radius: 8px 8px 0 0 !important;
        padding: 8px 16px !important;
    }
    .stTabs [aria-selected="true"] {
        background-color: #4a8af4 !important;
        color: white !important;
    }
    .stDataFrame {
        border-radius: 8px;
    }
    .stAlert {
        border-radius: 8px;
    }
</style>
""", unsafe_allow_html=True)

# Metrics
with st.container():
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown(f'<div class="metric-card"><h3>🔍 Searches</h3><h1>{len(history)}</h1></div>', unsafe_allow_html=True)
    with col2:
        pos = sum(1 for f in feedback if str(f.get("helpful", f.get("feedback", ""))).startswith("👍"))
        neg = sum(1 for f in feedback if str(f.get("helpful", f.get("feedback", ""))).startswith("👎"))
        st.markdown(f'<div class="metric-card"><h3>👍 / 👎 Feedback</h3><h1>{pos} / {neg}</h1></div>', unsafe_allow_html=True)
    with col3:
        st.markdown(f'<div class="metric-card"><h3>📁 Uploads</h3><h1>{len(uploads)}</h1></div>', unsafe_allow_html=True)

st.markdown("---")

# ── Upload Section ──────────────────────────────────────────────────
with st.expander("📂 Upload New Knowledge", expanded=False):
    u_files = st.file_uploader(
        "Drag & drop files here (PDF, TXT, CSV, Excel, JSON, DOCX, PPTX)",
        accept_multiple_files=True,
        type=["pdf", "txt", "csv", "xlsx", "xls", "json", "docx", "pptx"],
        help="Files will be processed and added to the vector database"
    )

    def extract_text(file) -> str | list[dict] | None:
        import fitz, docx
        from pptx import Presentation
        from PIL import Image
        import pytesseract, io

        ext = file.name.split(".")[-1].lower()
        
        try:
            if ext == "pdf":
                doc = fitz.open(stream=file.read(), filetype="pdf")
                texts = []
                for page in doc:
                    text = page.get_text("text").strip()
                    if not text:
                        pix = page.get_pixmap(dpi=300)
                        img = Image.open(io.BytesIO(pix.tobytes("png")))
                        text = pytesseract.image_to_string(img)
                    texts.append(text)
                return "\n".join(texts)
            
            elif ext == "txt":
                return file.read().decode("utf-8", "ignore")

            elif ext == "csv":
                return pd.read_csv(file).astype(str).to_csv(index=False)

            elif ext in ("xlsx", "xls"):
                sheets = pd.read_excel(file, sheet_name=None)
                return "\n".join(s.astype(str).to_csv(index=False) for s in sheets.values())

            elif ext == "docx":
                d = docx.Document(file)
                return "\n".join(p.text for p in d.paragraphs if p.text.strip())

            elif ext == "pptx":
                ppt = Presentation(file)
                return "\n".join(shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text"))

            elif ext == "json":
                return json.load(file)

        except Exception as e:
            st.error(f"❌ Error processing {file.name}: {str(e)[:200]}...")
            return None

    def chunk_text(txt: str, size=1200) -> list[str]:
        txt = re.sub(r"\s+", " ", txt.strip())
        return [txt[i:i+size] for i in range(0, len(txt), size)]

    if u_files and st.button("🚀 Push to Vector DB", type="primary"):
        progress_bar = st.progress(0)
        status_text = st.empty()
        new_rec, faq_entries = [], []
        
        for i, uf in enumerate(u_files):
            progress = (i + 1) / len(u_files)
            progress_bar.progress(progress)
            status_text.text(f"Processing {i+1}/{len(u_files)}: {uf.name}")
            
            with st.spinner(f"Extracting text from {uf.name}"):
                raw = extract_text(uf)
            
            if not raw:
                st.warning(f"⚠️ No text extracted from {uf.name}")
                continue

            if isinstance(raw, list) and all("text" in r for r in raw):
                final_chunks, ids, metas = [], [], []
                existing_ids = set(collection.get()["ids"])

                for r in raw:
                    text = str(r["text"]).strip()
                    if len(text) < 30:
                        continue
                    cid = r.get("id", f"ul_{uuid.uuid4().hex[:8]}")
                    if cid in existing_ids:
                        continue
                    meta = r.get("metadata", {}) or {}
                    meta.update({
                        "filename": uf.name,
                        "uploaded_at": now_iso()
                    })
                    ids.append(cid)
                    final_chunks.append(text)
                    metas.append(meta)
                    faq_entries.append({"text": text, "metadata": meta})

                if final_chunks:
                    embeds = embedder.encode(final_chunks).tolist()
                    collection.add(documents=final_chunks, ids=ids, embeddings=embeds, metadatas=metas)
                    new_rec.append({
                        "timestamp": now_iso(),
                        "filename": uf.name,
                        "size_kb": round(len(uf.getvalue()) / 1024, 2),
                        "chunks": len(final_chunks),
                        "deleted_at": ""
                    })
                    st.success(f"✅ {uf.name}: {len(final_chunks)} structured chunks inserted")
                continue

            chunks = chunk_text(str(raw))
            if not chunks:
                st.warning(f"⚠️ {uf.name}: no extractable text")
                continue
            
            existing_ids = set(collection.get()["ids"])
            ids, final_chunks = [], []

            for i, chunk in enumerate(chunks):
                cid = f"ul_{uuid.uuid4().hex[:8]}_{i}"
                if cid not in existing_ids:
                    ids.append(cid)
                    final_chunks.append(chunk)

            if not final_chunks:
                continue

            embeds = embedder.encode(final_chunks).tolist()
            metas = [{
                "filename": uf.name,
                "uploaded_at": now_iso(),
                "category": "PDF Upload",
                "recency": datetime.now().year,
                "priority": 5
            }] * len(final_chunks)

            collection.add(documents=final_chunks, ids=ids, embeddings=embeds, metadatas=metas)
            new_rec.append({
                "timestamp": now_iso(),
                "filename": uf.name,
                "size_kb": round(len(uf.getvalue()) / 1024, 2),
                "chunks": len(final_chunks),
                "deleted_at": ""
            })
            faq_entries += [{"text": chunk, "metadata": meta} for chunk, meta in zip(final_chunks, metas)]

        progress_bar.empty()
        status_text.empty()
        
        if new_rec:
            uploads.extend(new_rec)
            save_json(UPLOAD_LOG_FILE, uploads)
            if faq_entries:
                append_to_faq_json(faq_entries)
            st.success("✅ All files processed successfully!")
            st.balloons()
            st.rerun()

# ── Chunks Table ────────────────────────────────────────────────────
st.subheader("📚 Knowledge Base Contents")
db = collection.get(include=["documents", "metadatas"])
if not db.get("ids"):
    st.info("Vector DB is empty. Upload files to add content.")
    db_df = pd.DataFrame()  # Define empty db_df to avoid NameError
else:
    db_df = pd.DataFrame({
        "chunk_id": db["ids"],
        "filename": [m.get("filename", "faq") for m in db["metadatas"]],
        "uploaded": [m.get("uploaded_at", "-") for m in db["metadatas"]],
        "category": [m.get("category", "faq") for m in db["metadatas"]],
        "preview": [d[:100] + ("…" if len(d) > 100 else "") for d in db["documents"]],
    })


    # Add search and filter functionality
    with st.expander("🔍 Search & Filter", expanded=False):
        search_col1, search_col2 = st.columns(2)
        with search_col1:
            search_text = st.text_input("Search content")
        with search_col2:
            filter_category = st.selectbox("Filter by category", ["All"] + sorted(db_df["category"].unique()))
        
        if search_text:
            db_df = db_df[db_df["preview"].str.contains(search_text, case=False, na=False)]
        if filter_category != "All":
            db_df = db_df[db_df["category"] == filter_category]

    st.dataframe(
        db_df,
        use_container_width=True,
        column_config={
            "chunk_id": "ID",
            "filename": "File",
            "uploaded": "Upload Date",
            "category": "Category",
            "preview": "Content Preview"
        },
        hide_index=True
    )

# ── Deletion Tabs ───────────────────────────────────────────────────
tab1, tab2, tab3 = st.tabs(["🗂 Delete by File", "🔍 Delete by Chunk ID", "🧨 Danger Zone"])

with tab1:
    st.subheader("Delete by File")
    filenames = sorted(db_df["filename"].unique()) if not db_df.empty else []
    if filenames:
        selected_file = st.selectbox("Select file to delete", filenames)
        if st.button("❌ Delete Selected File", type="secondary"):
            ids_to_del = db_df.loc[db_df["filename"] == selected_file, "chunk_id"].tolist()
            if ids_to_del:
                collection.delete(ids=ids_to_del)
                for rec in uploads:
                    if rec.get("filename") == selected_file and not rec.get("deleted_at", ""):
                        rec["deleted_at"] = now_iso()
                save_json(UPLOAD_LOG_FILE, uploads)
                st.success(f"✅ Deleted {len(ids_to_del)} chunks from '{selected_file}'")
                st.rerun()
    else:
        st.info("No files available for deletion")

with tab2:
    st.subheader("Delete by Chunk ID")
    if not db_df.empty:
        selected_ids = st.multiselect("Select chunk IDs to delete", db_df["chunk_id"].tolist())
        if selected_ids and st.button("❌ Delete Selected Chunks", type="secondary"):
            collection.delete(ids=selected_ids)
            st.success(f"✅ Deleted {len(selected_ids)} chunk(s)")
            st.rerun()
    else:
        st.info("No chunks available for deletion")

with tab3:
    st.subheader("Danger Zone")
    st.warning("These actions are irreversible!")
    
    if not db_df.empty:
        if st.button("❌ Delete ALL Data from Vector DB", type="primary"):
            with st.spinner("Deleting all chunks..."):
                all_ids = db_df["chunk_id"].tolist()
                collection.delete(ids=all_ids)
                for rec in uploads:
                    if not rec.get("deleted_at"):
                        rec["deleted_at"] = now_iso()
                save_json(UPLOAD_LOG_FILE, uploads)
            st.success(f"✅ Deleted all {len(all_ids)} chunks")
            st.rerun()
    else:
        st.info("Vector DB is already empty")

# ── Logs Overview ───────────────────────────────────────────────────
st.header("📊 Analytics & Logs")

def safe_json_load(obj):
    if isinstance(obj, str):
        try: return json.loads(obj)
        except Exception: return obj
    return obj

def first_meta(ctxs, key):
    ctxs = safe_json_load(ctxs)
    if isinstance(ctxs, list) and ctxs:
        for c in ctxs:
            if isinstance(c, dict):
                if key in c: return c[key]
                if isinstance(c.get("meta"), dict):
                    return c["meta"].get(key, "")
    return ""

def ensure_cols(df: pd.DataFrame, columns: list[str]) -> pd.DataFrame:
    for col in columns:
        if col not in df.columns:
            df[col] = ""
    return df[columns]

# 🔍 Search History
with st.expander("🔍 Search History", expanded=True):
    df_search = pd.DataFrame(history)

    if not df_search.empty:
        df_search["timestamp"] = pd.to_datetime(df_search["timestamp"], utc=True, errors="coerce")
        df_search["timestamp"] = df_search["timestamp"].dt.tz_convert("Asia/Kolkata").dt.strftime("%Y-%m-%d %H:%M:%S")
        df_search = df_search.sort_values("timestamp", ascending=False)

        if "top_contexts" in df_search.columns:
            df_search["source_file"] = df_search["top_contexts"].apply(lambda c: first_meta(c, "filename"))
            df_search["category"]    = df_search["top_contexts"].apply(lambda c: first_meta(c, "category"))
            df_search["chunk_id"]    = df_search["top_contexts"].apply(lambda c: first_meta(c, "chunk_id"))
        df_search = df_search.drop(columns=["top_contexts"], errors="ignore")
        df_search = ensure_cols(df_search, ["timestamp","user_query","ai_answer","source_file","category","chunk_id"])

        # Enhanced filtering
        with st.expander("🔎 Advanced Filters", expanded=False):
            col1, col2, col3 = st.columns(3)
            with col1:
                date_filter = st.date_input("Filter by date")
            with col2:
                category_filter = st.selectbox("Filter by category", ["All"] + sorted(df_search["category"].unique().tolist()))
            with col3:
                query_filter = st.text_input("Filter by query or answer")

        if date_filter:
    # Convert to datetime if not already
             df_search["timestamp"] = pd.to_datetime(df_search["timestamp"], errors='coerce')
    # Filter by date
             df_search = df_search[df_search["timestamp"].dt.date == date_filter]
        if category_filter != "All":
            df_search = df_search[df_search["category"] == category_filter]
        if query_filter:
            df_search = df_search[
                df_search["user_query"].str.contains(query_filter, case=False, na=False) |
                df_search["ai_answer"].str.contains(query_filter, case=False, na=False)
            ]

        st.dataframe(
            df_search,
            use_container_width=True,
            column_config={
                "timestamp": "Time",
                "user_query": "Query",
                "ai_answer": "Answer",
                "source_file": "Source",
                "category": "Category",
                "chunk_id": "Chunk ID"
            },
            hide_index=True
        )

        col1, col2 = st.columns(2)
        with col1:
            st.download_button("📥 Download JSON", json.dumps(history, indent=2), file_name="search_history.json")
        with col2:
            st.download_button("📄 Download CSV", df_search.to_csv(index=False), file_name="search_history.csv")
    else:
        st.info("No search history yet.")

# 🧠 LLM Queries
with st.expander("🧠 LLM Queries", expanded=False):
    df_llm = to_df(query_log)
    if not df_llm.empty:
        if "top_contexts" in df_llm.columns:
            df_llm["source_file"] = df_llm["top_contexts"].apply(lambda c: first_meta(c, "filename"))
        df_llm = df_llm.drop(columns=["top_contexts"], errors="ignore")
        df_llm = ensure_cols(df_llm, ["timestamp","user_query","llm_prompt","llm_response","source_file"])
        
        st.dataframe(
            df_llm,
            use_container_width=True,
            column_config={
                "timestamp": "Time",
                "user_query": "User Query",
                "llm_prompt": "LLM Prompt",
                "llm_response": "LLM Response",
                "source_file": "Source File"
            },
            hide_index=True
        )
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button("📥 Download JSON", json.dumps(query_log, indent=2), file_name="llm_queries.json")
        with col2:
            st.download_button("📄 Download CSV", df_llm.to_csv(index=False), file_name="llm_queries.csv")
    else:
        st.info("No LLM queries logged yet.")

# 💬 Feedback
with st.expander("💬 User Feedback", expanded=False):
    df_fb = to_df(feedback)
    if not df_fb.empty:
        if "helpful" in df_fb.columns and "feedback" not in df_fb.columns:
            df_fb = df_fb.rename(columns={"helpful": "feedback"})
        df_fb["feedback"] = df_fb["feedback"].astype(str)
        df_fb = ensure_cols(df_fb, ["timestamp","user_query","ai_answer","feedback","comment"])
        
        # Feedback analysis
        pos_fb = df_fb[df_fb["feedback"].str.startswith("👍")]
        neg_fb = df_fb[df_fb["feedback"].str.startswith("👎")]
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Positive Feedback", len(pos_fb))
        with col2:
            st.metric("Negative Feedback", len(neg_fb))
        
        st.dataframe(
            df_fb,
            use_container_width=True,
            column_config={
                "timestamp": "Time",
                "user_query": "User Query",
                "ai_answer": "AI Answer",
                "feedback": "Feedback",
                "comment": "Comment"
            },
            hide_index=True
        )
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button("📥 Download JSON", json.dumps(feedback, indent=2), file_name="feedback.json")
        with col2:
            st.download_button("📄 Download CSV", df_fb.to_csv(index=False), file_name="feedback.csv")
    else:
        st.info("No feedback received yet.")

# 📜 Uploads
with st.expander("📜 File Upload History", expanded=False):
    df_up = to_df(uploads)
    if not df_up.empty:
        df_up = ensure_cols(df_up, ["timestamp","filename","size_kb","chunks","deleted_at"])
        
        # Upload stats
        active_uploads = df_up[df_up["deleted_at"] == ""]
        deleted_uploads = df_up[df_up["deleted_at"] != ""]
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Active Uploads", len(active_uploads))
        with col2:
            st.metric("Deleted Uploads", len(deleted_uploads))
        
        st.dataframe(
            df_up,
            use_container_width=True,
            column_config={
                "timestamp": "Time",
                "filename": "Filename",
                "size_kb": "Size (KB)",
                "chunks": "Chunks",
                "deleted_at": "Deleted At"
            },
            hide_index=True
        )
        
        col1, col2 = st.columns(2)
        with col1:
            st.download_button("📥 Download JSON", json.dumps(uploads, indent=2), file_name="uploads.json")
        with col2:
            st.download_button("📄 Download CSV", df_up.to_csv(index=False), file_name="uploads.csv")
    else:
        st.info("No files uploaded yet.")

# 📊 Analytics
with st.expander("📊 Visual Analytics", expanded=False):
    tab1, tab2 = st.tabs(["Search Activity", "Feedback Analysis"])
    
    with tab1:
        if history:
            df_h = to_df(history)
            df_h["date"] = df_h["timestamp"].dt.date
            df_h["hour"] = df_h["timestamp"].dt.hour
            
            col1, col2 = st.columns(2)
            with col1:
                st.altair_chart(
                    alt.Chart(df_h.groupby("date").size().reset_index(name="count"))
                       .mark_bar()
                       .encode(
                           x="date:T",
                           y="count:Q",
                           tooltip=["date:T","count:Q"],
                           color=alt.Color("count:Q", scale=alt.Scale(scheme="blues"))
                       )
                       .properties(title="Daily Search Volume", width="container"),
                    use_container_width=True
                )
            with col2:
                st.altair_chart(
                    alt.Chart(df_h.groupby("hour").size().reset_index(name="count"))
                       .mark_line(point=True)
                       .encode(
                           x="hour:O",
                           y="count:Q",
                           tooltip=["hour:O","count:Q"]
                       )
                       .properties(title="Hourly Search Pattern", width="container"),
                    use_container_width=True
                )
        else:
            st.info("No search history for analytics")
    
    with tab2:
        if feedback:
            df_fb = to_df(feedback)
            if "helpful" in df_fb.columns and "feedback" not in df_fb.columns:
                df_fb = df_fb.rename(columns={"helpful": "feedback"})
            df_fb["vote"] = df_fb["feedback"].astype(str).str[:2]
            df_fb["date"] = pd.to_datetime(df_fb["timestamp"]).dt.date
            
            col1, col2 = st.columns(2)
            with col1:
                st.altair_chart(
                    alt.Chart(df_fb)
                       .mark_arc()
                       .encode(
                           theta="count():Q",
                           color=alt.Color("vote:N", legend=None,
                                           scale=alt.Scale(range=["#841306","#5bee60"])),
                           tooltip=["vote:N","count():Q"]
                       )
                       .properties(title="Feedback Ratio", width="container"),
                    use_container_width=True
                )
            with col2:
                st.altair_chart(
                    alt.Chart(df_fb.groupby("date")["vote"].value_counts().reset_index(name="count"))
                       .mark_area()
                       .encode(
                           x="date:T",
                           y="count:Q",
                           color="vote:N",
                           tooltip=["date:T","vote:N","count:Q"]
                       )
                       .properties(title="Feedback Trend", width="container"),
                    use_container_width=True
                )
        else:
            st.info("No feedback data for analytics")

# 📝 FAQ Manager
with st.expander("📝 FAQ Management", expanded=False):
    st.subheader("FAQ Entries")
    
    # Add new FAQ
    with st.form("add_faq"):
        st.write("### Add New FAQ Entry")
        col1, col2 = st.columns(2)
        text = st.text_area("FAQ Text (Answer)", height=100, help="The detailed answer to the FAQ")
        category = col1.text_input("Category", help="e.g., account, orders, payment")
        recency = col2.number_input("Recency Year", min_value=2000, max_value=2100, value=datetime.now().year)
        priority = col1.slider("Priority (1-10)", min_value=1, max_value=10, value=5, help="Higher priority items appear first in results")
        source = col2.text_input("Source", value="manual_entry", help="Source of this FAQ entry")
        
        if st.form_submit_button("✅ Add FAQ Entry", type="primary"):
            cleaned, tag = clean_and_tag(text)
            if tag != "informative":
                st.warning("Text is not informative enough. Please provide more detailed content.")
            else:
                faq_data = load_json(FAQ_JSON_FILE)
                max_id = max((int(e.get("id", "0")) for e in faq_data), default=0)
                next_id = str(max_id + 1)

                entry = {
                    "id": next_id,
                    "text": cleaned,
                    "metadata": {
                        "category": category.strip().lower(),
                        "source": source.strip(),
                        "recency": int(recency),
                        "priority": int(priority),
                        "uploaded_at": now_iso()
                    }
                }

                # Save to JSON
                faq_data.append(entry)
                save_json(FAQ_JSON_FILE, faq_data)
                st.success(f"Added FAQ entry with ID: {next_id}")

                # Save to Vector DB
                try:
                    vector_id = f"faq_{next_id}"
                    existing_ids = set(collection.get()["ids"])
                    if vector_id not in existing_ids:
                        embedding = embedder.encode([cleaned])[0].tolist()
                        collection.add(
                            documents=[cleaned],
                            metadatas=[entry["metadata"]],
                            ids=[vector_id],
                            embeddings=[embedding]
                        )
                        st.success(f"Also added to VectorDB as `{vector_id}`")
                    else:
                        st.warning("This ID already exists in VectorDB")
                except Exception as e:
                    st.error(f"VectorDB error: {str(e)[:200]}...")
                st.rerun()

    # View existing FAQs
    st.write("### Existing FAQ Entries")
    faq_data = load_json(FAQ_JSON_FILE)
    if not faq_data:
        st.info("No FAQ entries yet.")
    else:
        df_faq = pd.DataFrame(faq_data)
        df_faq["category"] = df_faq["metadata"].apply(lambda m: m.get("category", ""))
        df_faq["priority"] = df_faq["metadata"].apply(lambda m: m.get("priority", ""))
        df_faq["recency"] = df_faq["metadata"].apply(lambda m: m.get("recency", ""))
        df_faq["source"] = df_faq["metadata"].apply(lambda m: m.get("source", ""))
        df_faq = df_faq.drop(columns=["metadata"])

        # Search and filter
        with st.expander("🔍 Search FAQs", expanded=False):
            search_col1, search_col2 = st.columns(2)
            with search_col1:
                faq_search = st.text_input("Search FAQ text")
            with search_col2:
                faq_category = st.selectbox("Filter by category", ["All"] + sorted(df_faq["category"].unique().tolist()))
            
            if faq_search:
                df_faq = df_faq[df_faq["text"].str.contains(faq_search, case=False, na=False)]
            if faq_category != "All":
                df_faq = df_faq[df_faq["category"] == faq_category]

        st.dataframe(
            df_faq,
            use_container_width=True,
            column_config={
                "id": "ID",
                "text": "FAQ Content",
                "category": "Category",
                "priority": "Priority",
                "recency": "Year",
                "source": "Source"
            },
            hide_index=True
        )

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; padding: 20px;">
    <p>Made with ❤️ by Joyjeet Roy</p>
    <p>Version 2.0 • Updated {}</p>
</div>
""".format(datetime.now().strftime("%Y-%m-%d")), unsafe_allow_html=True)